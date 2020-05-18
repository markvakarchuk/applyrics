from pptx import Presentation
from testsongs import *
import os
from pathlib import Path
import json
from json import JSONEncoder
import codecs
# coding: latin-1
class song:
    def __init__(self, title, versus):
        self.title = title
        self. versus = versus

    def print(self):
        print(str(self.title) + "\n")
        print(str(self.versus) + "\n")

    def toJson(self):
        return self.__dict__
# class song:
#     title = ""
#     versus = []




directory = "testsongs"
files_in_path_ext = os.listdir(directory)
files_in_path_noext = []

# for files in os.listdir(directory):
#     f = os.path.splitext(files)[0]
#     # print(f)
#     files_in_path_noext.append(f)

# print(files_in_path_noext)
data = []
i = 0
s = ""
for file in files_in_path_ext:
    s = "s" + str(i)
    versus = []
    title = str(os.path.splitext(file)[0])
    print(">" + title)
    filepath = os.path.join(directory, file)
    # print(filepath)
    # print(f)
    prs = Presentation(filepath)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # print(shape.text)
                versus.append(shape.text)
    # print(versus)
    s = song(title, versus)
    # s.title = title
    # s.versus = versus
    # songs.append(s)
    # print("************************\n")
    # i+=1
    # s.print
    x = s.toJson()
    data.append(x)


# print(data)

with codecs.open('data_test.json', 'a+', 'utf-8') as f:
    json.dump(data, f, ensure_ascii=False )